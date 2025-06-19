import os
import re
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
import PyPDF2
from docx import Document
from pptx import Presentation
from config.settings import Settings

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        self.supported_extensions = Settings.SUPPORTED_EXTENSIONS
        self.chunk_size = Settings.CHUNK_SIZE
        self.chunk_overlap = Settings.CHUNK_OVERLAP
    
    def extract_text_from_file(self, file_path: Path) -> str:
        """Extract text content from various file formats"""
        try:
            extension = file_path.suffix.lower()
            
            if extension == '.pdf':
                return self._extract_pdf_text(file_path)
            elif extension == '.docx':
                return self._extract_docx_text(file_path)
            elif extension == '.pptx':
                return self._extract_pptx_text(file_path)
            elif extension in ['.txt', '.md']:
                return self._extract_text_file(file_path)
            else:
                logger.warning(f"Unsupported file type: {extension}")
                return ""
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = ""
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                    if row_text:
                        text += row_text + "\n"
            
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""
    
    def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {slide_num + 1} ---\n"
                has_content = False
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                        has_content = True
                
                if has_content:
                    text += slide_text
            
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def _extract_text_file(self, file_path: Path) -> str:
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error reading text file {file_path}: {e}")
                return ""
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return ""
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize text content"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep basic punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\{\}]', ' ', text)
        
        # Normalize line breaks
        text = re.sub(r'\n+', '\n', text)
        
        # Remove leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def chunk_text(self, text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[str]:
        """Split text into overlapping chunks"""
        if not text:
            return []
        
        chunk_size = chunk_size or self.chunk_size
        chunk_overlap = chunk_overlap or self.chunk_overlap
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # If this isn't the last chunk, try to break at a sentence boundary
            if end < len(text):
                # Look for sentence endings within the last 100 characters of the chunk
                search_start = max(start, end - 100)
                sentence_end = text.rfind('.', search_start, end)
                if sentence_end > start:
                    end = sentence_end + 1
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position for next chunk
            start = end - chunk_overlap
            if start >= len(text):
                break
        
        return chunks
    
    def get_file_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from file"""
        try:
            stat = file_path.stat()
            return {
                "file_path": str(file_path),
                "file_name": file_path.name,
                "file_size": stat.st_size,
                "file_extension": file_path.suffix.lower(),
                "created_time": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified_time": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "accessed_time": datetime.fromtimestamp(stat.st_atime).isoformat(),
            }
        except Exception as e:
            logger.error(f"Error getting file metadata for {file_path}: {e}")
            return {}
    
    def process_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Process a single file and return chunks with metadata"""
        try:
            # Check if file is supported
            if file_path.suffix.lower() not in self.supported_extensions:
                logger.warning(f"Unsupported file type: {file_path.suffix}")
                return []
            
            # Extract text
            raw_text = self.extract_text_from_file(file_path)
            if not raw_text:
                logger.warning(f"No text extracted from {file_path}")
                return []
            
            # Clean text
            cleaned_text = self.clean_text(raw_text)
            if not cleaned_text:
                logger.warning(f"No content after cleaning for {file_path}")
                return []
            
            # Chunk text
            chunks = self.chunk_text(cleaned_text)
            if not chunks:
                logger.warning(f"No chunks created for {file_path}")
                return []
            
            # Get file metadata
            metadata = self.get_file_metadata(file_path)
            
            # Create chunk documents
            documents = []
            for i, chunk in enumerate(chunks):
                chunk_metadata = metadata.copy()
                chunk_metadata.update({
                    "chunk_id": i,
                    "total_chunks": len(chunks),
                    "chunk_size": len(chunk),
                    "chunk_text": chunk[:100] + "..." if len(chunk) > 100 else chunk,  # Preview
                })
                
                documents.append({
                    "content": chunk,
                    "metadata": chunk_metadata
                })
            
            logger.info(f"Processed {file_path}: {len(chunks)} chunks created")
            return documents
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return []
    
    def process_directory(self, directory_path: Path, recursive: bool = True) -> List[Dict[str, Any]]:
        """Process all supported files in a directory"""
        all_documents = []
        
        try:
            if recursive:
                file_paths = list(directory_path.rglob("*"))
            else:
                file_paths = list(directory_path.iterdir())
            
            # Filter for files with supported extensions
            supported_files = [
                path for path in file_paths 
                if path.is_file() and path.suffix.lower() in self.supported_extensions
            ]
            
            logger.info(f"Found {len(supported_files)} supported files in {directory_path}")
            
            for file_path in supported_files:
                documents = self.process_file(file_path)
                all_documents.extend(documents)
            
            logger.info(f"Total documents processed: {len(all_documents)}")
            return all_documents
            
        except Exception as e:
            logger.error(f"Error processing directory {directory_path}: {e}")
            return all_documents 