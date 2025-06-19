import asyncio
import json
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence
from mcp.server import Server
from mcp.server.models import InitializationOptions
from mcp.server.stdio import stdio_server
from mcp.types import (
    CallToolRequest,
    CallToolResult,
    ListToolsRequest,
    ListToolsResult,
    Tool,
    TextContent,
    ImageContent,
    EmbeddedResource,
    LoggingLevel,
)
import PyPDF2
from docx import Document
from pptx import Presentation
import os
import mimetypes
from datetime import datetime

logger = logging.getLogger(__name__)

class FileSystemServer:
    def __init__(self, root_path: str = None):
        self.root_path = Path(root_path) if root_path else Path.home()
        self.server = Server("file-system-server")
        self.setup_handlers()
    
    def setup_handlers(self):
        @self.server.list_tools()
        async def handle_list_tools() -> ListToolsResult:
            return ListToolsResult(
                tools=[
                    Tool(
                        name="list_files",
                        description="List files and directories in a given path",
                        inputSchema={
                            "type": "object",
                            "properties": {
                                "path": {
                                    "type": "string",
                                    "description": "Path to list files from (defaults to root)"
                                }
                            }
                        }
                    ),
                    Tool(
                        name="read_file",
                        description="Read the content of a file",
                        inputSchema={
                            "type": "object",
                            "properties": {
                                "path": {
                                    "type": "string",
                                    "description": "Path to the file to read"
                                },
                                "encoding": {
                                    "type": "string",
                                    "description": "File encoding (default: utf-8)"
                                }
                            },
                            "required": ["path"]
                        }
                    ),
                    Tool(
                        name="get_file_metadata",
                        description="Get metadata about a file or directory",
                        inputSchema={
                            "type": "object",
                            "properties": {
                                "path": {
                                    "type": "string",
                                    "description": "Path to the file or directory"
                                }
                            },
                            "required": ["path"]
                        }
                    ),
                    Tool(
                        name="extract_document_content",
                        description="Extract text content from various document formats (PDF, DOCX, PPTX, TXT, MD)",
                        inputSchema={
                            "type": "object",
                            "properties": {
                                "path": {
                                    "type": "string",
                                    "description": "Path to the document file"
                                }
                            },
                            "required": ["path"]
                        }
                    ),
                    Tool(
                        name="search_files",
                        description="Search for files by name or content",
                        inputSchema={
                            "type": "object",
                            "properties": {
                                "query": {
                                    "type": "string",
                                    "description": "Search query"
                                },
                                "path": {
                                    "type": "string",
                                    "description": "Root path to search in (defaults to root)"
                                },
                                "file_types": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": "File extensions to search (e.g., ['.txt', '.pdf'])"
                                }
                            },
                            "required": ["query"]
                        }
                    )
                ]
            )

        @self.server.call_tool()
        async def handle_call_tool(name: str, arguments: Dict[str, Any]) -> CallToolResult:
            try:
                if name == "list_files":
                    return await self._list_files(arguments)
                elif name == "read_file":
                    return await self._read_file(arguments)
                elif name == "get_file_metadata":
                    return await self._get_file_metadata(arguments)
                elif name == "extract_document_content":
                    return await self._extract_document_content(arguments)
                elif name == "search_files":
                    return await self._search_files(arguments)
                else:
                    return CallToolResult(
                        content=[TextContent(type="text", text=f"Unknown tool: {name}")]
                    )
            except Exception as e:
                logger.error(f"Error in tool {name}: {e}")
                return CallToolResult(
                    content=[TextContent(type="text", text=f"Error: {str(e)}")]
                )

    async def _list_files(self, arguments: Dict[str, Any]) -> CallToolResult:
        path = arguments.get("path", str(self.root_path))
        target_path = Path(path)
        
        if not target_path.exists():
            return CallToolResult(
                content=[TextContent(type="text", text=f"Path does not exist: {path}")]
            )
        
        try:
            items = []
            for item in target_path.iterdir():
                item_info = {
                    "name": item.name,
                    "path": str(item),
                    "is_file": item.is_file(),
                    "is_dir": item.is_dir(),
                    "size": item.stat().st_size if item.is_file() else None,
                    "modified": datetime.fromtimestamp(item.stat().st_mtime).isoformat()
                }
                items.append(item_info)
            
            # Sort: directories first, then files
            items.sort(key=lambda x: (not x["is_dir"], x["name"].lower()))
            
            return CallToolResult(
                content=[TextContent(
                    type="text", 
                    text=f"Files and directories in {path}:\n\n" + 
                         "\n".join([
                             f"{'📁' if item['is_dir'] else '📄'} {item['name']} "
                             f"({item['size']} bytes)" if item['is_file'] else f"{item['name']}"
                             for item in items
                         ])
                )]
            )
        except Exception as e:
            return CallToolResult(
                content=[TextContent(type="text", text=f"Error listing files: {str(e)}")]
            )

    async def _read_file(self, arguments: Dict[str, Any]) -> CallToolResult:
        path = arguments["path"]
        encoding = arguments.get("encoding", "utf-8")
        target_path = Path(path)
        
        if not target_path.exists():
            return CallToolResult(
                content=[TextContent(type="text", text=f"File does not exist: {path}")]
            )
        
        if not target_path.is_file():
            return CallToolResult(
                content=[TextContent(type="text", text=f"Path is not a file: {path}")]
            )
        
        try:
            with open(target_path, 'r', encoding=encoding) as f:
                content = f.read()
            
            return CallToolResult(
                content=[TextContent(type="text", text=content)]
            )
        except Exception as e:
            return CallToolResult(
                content=[TextContent(type="text", text=f"Error reading file: {str(e)}")]
            )

    async def _get_file_metadata(self, arguments: Dict[str, Any]) -> CallToolResult:
        path = arguments["path"]
        target_path = Path(path)
        
        if not target_path.exists():
            return CallToolResult(
                content=[TextContent(type="text", text=f"Path does not exist: {path}")]
            )
        
        try:
            stat = target_path.stat()
            mime_type, _ = mimetypes.guess_type(str(target_path))
            
            metadata = {
                "name": target_path.name,
                "path": str(target_path),
                "is_file": target_path.is_file(),
                "is_dir": target_path.is_dir(),
                "size": stat.st_size,
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "accessed": datetime.fromtimestamp(stat.st_atime).isoformat(),
                "mime_type": mime_type,
                "extension": target_path.suffix.lower()
            }
            
            return CallToolResult(
                content=[TextContent(
                    type="text", 
                    text=f"Metadata for {path}:\n\n" + 
                         "\n".join([f"{k}: {v}" for k, v in metadata.items()])
                )]
            )
        except Exception as e:
            return CallToolResult(
                content=[TextContent(type="text", text=f"Error getting metadata: {str(e)}")]
            )

    async def _extract_document_content(self, arguments: Dict[str, Any]) -> CallToolResult:
        path = arguments["path"]
        target_path = Path(path)
        
        if not target_path.exists():
            return CallToolResult(
                content=[TextContent(type="text", text=f"File does not exist: {path}")]
            )
        
        if not target_path.is_file():
            return CallToolResult(
                content=[TextContent(type="text", text=f"Path is not a file: {path}")]
            )
        
        try:
            content = ""
            extension = target_path.suffix.lower()
            
            if extension == '.pdf':
                content = self._extract_pdf_content(target_path)
            elif extension == '.docx':
                content = self._extract_docx_content(target_path)
            elif extension == '.pptx':
                content = self._extract_pptx_content(target_path)
            elif extension in ['.txt', '.md']:
                with open(target_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            else:
                return CallToolResult(
                    content=[TextContent(type="text", text=f"Unsupported file type: {extension}")]
                )
            
            return CallToolResult(
                content=[TextContent(type="text", text=content)]
            )
        except Exception as e:
            return CallToolResult(
                content=[TextContent(type="text", text=f"Error extracting content: {str(e)}")]
            )

    def _extract_pdf_content(self, file_path: Path) -> str:
        """Extract text content from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PDF content: {str(e)}")

    def _extract_docx_content(self, file_path: Path) -> str:
        """Extract text content from DOCX file"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting DOCX content: {str(e)}")

    def _extract_pptx_content(self, file_path: Path) -> str:
        """Extract text content from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PPTX content: {str(e)}")

    async def _search_files(self, arguments: Dict[str, Any]) -> CallToolResult:
        query = arguments["query"].lower()
        search_path = arguments.get("path", str(self.root_path))
        file_types = arguments.get("file_types", [])
        
        target_path = Path(search_path)
        if not target_path.exists():
            return CallToolResult(
                content=[TextContent(type="text", text=f"Search path does not exist: {search_path}")]
            )
        
        try:
            results = []
            for root, dirs, files in os.walk(target_path):
                for file in files:
                    file_path = Path(root) / file
                    
                    # Filter by file type if specified
                    if file_types and file_path.suffix.lower() not in file_types:
                        continue
                    
                    # Search in filename
                    if query in file.lower():
                        results.append(str(file_path))
                        continue
                    
                    # Search in content for supported file types
                    if file_path.suffix.lower() in ['.txt', '.md', '.pdf', '.docx', '.pptx']:
                        try:
                            content = await self._extract_document_content({"path": str(file_path)})
                            if query in content.content[0].text.lower():
                                results.append(str(file_path))
                        except:
                            continue
            
            if results:
                return CallToolResult(
                    content=[TextContent(
                        type="text", 
                        text=f"Found {len(results)} files matching '{query}':\n\n" + 
                             "\n".join([f"📄 {result}" for result in results[:20]])  # Limit to 20 results
                    )]
                )
            else:
                return CallToolResult(
                    content=[TextContent(type="text", text=f"No files found matching '{query}'")]
                )
        except Exception as e:
            return CallToolResult(
                content=[TextContent(type="text", text=f"Error searching files: {str(e)}")]
            )

    async def run(self):
        """Run the MCP server"""
        async with stdio_server() as (read_stream, write_stream):
            await self.server.run(
                read_stream,
                write_stream,
                InitializationOptions(
                    server_name="file-system-server",
                    server_version="1.0.0",
                    capabilities=self.server.get_capabilities(
                        notification_options=None,
                        experimental_capabilities=None,
                    ),
                ),
            )

def main():
    """Main entry point for the MCP file system server"""
    import argparse
    
    parser = argparse.ArgumentParser(description="MCP File System Server")
    parser.add_argument("--root-path", type=str, help="Root path for file operations")
    args = parser.parse_args()
    
    server = FileSystemServer(root_path=args.root_path)
    
    try:
        asyncio.run(server.run())
    except KeyboardInterrupt:
        logger.info("Server stopped by user")
    except Exception as e:
        logger.error(f"Server error: {e}")

if __name__ == "__main__":
    main()
