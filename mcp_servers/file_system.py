"""
FastMCP Semantic Filesystem Server
A semantic search-enabled filesystem server using FastMCP for multi-format file operations.
"""

import os
import sys
import pathspec
import fnmatch
from typing import List, Optional, Dict, Any, Union
from pathlib import Path
import mimetypes
from datetime import datetime
import hashlib
from dataclasses import dataclass
import json
import re

# Document processing imports
import PyPDF2
import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
import csv
import xml.etree.ElementTree as ET
from PIL import Image
import eyed3  # For audio metadata

# Semantic search imports
from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import pickle

from mcp.server.fastmcp import FastMCP

# Configuration constants
DEFAULT_IGNORE_PATTERNS = [
    ".git",
    "__pycache__",
    "*.pyc",
    ".venv",
    "venv",
    ".env",
    ".idea",
    ".vscode",
    "*.egg-info",
    "dist",
    "build",
    ".pytest_cache",
    ".coverage",
    "htmlcov",
    ".DS_Store",
    "Thumbs.db",
    "node_modules",
    "*.log",
    "*.tmp",
    "*.cache",
]

@dataclass
class FileMetadata:
    """Enhanced file metadata structure"""
    path: str
    size: int
    modified: datetime
    mime_type: str
    encoding: str
    hash: str
    content_type: str  # text, document, spreadsheet, presentation, etc.
    extractable: bool  # whether content can be extracted
    
@dataclass
class DocumentContent:
    """Extracted document content structure"""
    text: str
    metadata: Dict[str, Any]
    pages: Optional[List[str]] = None
    sheets: Optional[List[Dict[str, Any]]] = None
    slides: Optional[List[Dict[str, Any]]] = None

class DocumentExtractor:
    """Handles extraction of content from various document formats"""
    
    @staticmethod
    def extract_pdf(file_path: str) -> DocumentContent:
        """Extract content from PDF files"""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                pages = []
                
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    pages.append(page_text)
                    text += f"\n--- Page {page_num + 1} ---\n{page_text}"
                
                metadata = {
                    "num_pages": len(reader.pages),
                    "title": reader.metadata.get('/Title', '') if reader.metadata else '',
                    "author": reader.metadata.get('/Author', '') if reader.metadata else '',
                    "subject": reader.metadata.get('/Subject', '') if reader.metadata else '',
                }
                
                return DocumentContent(text=text, metadata=metadata, pages=pages)
        except Exception as e:
            return DocumentContent(text=f"Error extracting PDF: {str(e)}", metadata={})
    
    @staticmethod
    def extract_excel(file_path: str) -> DocumentContent:
        """Extract content from Excel files"""
        try:
            # Try pandas first for broader format support
            try:
                excel_data = pd.read_excel(file_path, sheet_name=None)
                text = ""
                sheets = []
                
                for sheet_name, df in excel_data.items():
                    sheet_text = f"\n--- Sheet: {sheet_name} ---\n"
                    sheet_text += df.to_string(index=False)
                    text += sheet_text
                    
                    sheets.append({
                        "name": sheet_name,
                        "rows": len(df),
                        "columns": len(df.columns),
                        "column_names": df.columns.tolist(),
                        "content": df.to_dict('records')[:100]  # Limit to first 100 rows
                    })
                
                metadata = {
                    "num_sheets": len(excel_data),
                    "sheet_names": list(excel_data.keys())
                }
                
                return DocumentContent(text=text, metadata=metadata, sheets=sheets)
                
            except Exception:
                # Fallback to openpyxl for xlsx files
                wb = load_workbook(file_path, read_only=True)
                text = ""
                sheets = []
                
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_text = f"\n--- Sheet: {sheet_name} ---\n"
                    sheet_data = []
                    
                    for row in ws.iter_rows(values_only=True):
                        if any(cell is not None for cell in row):
                            row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                            sheet_text += row_text + '\n'
                            sheet_data.append(row)
                    
                    text += sheet_text
                    sheets.append({
                        "name": sheet_name,
                        "rows": len(sheet_data),
                        "content": sheet_data[:100]  # Limit to first 100 rows
                    })
                
                metadata = {
                    "num_sheets": len(wb.sheetnames),
                    "sheet_names": wb.sheetnames
                }
                
                return DocumentContent(text=text, metadata=metadata, sheets=sheets)
                
        except Exception as e:
            return DocumentContent(text=f"Error extracting Excel: {str(e)}", metadata={})
    
    @staticmethod
    def extract_powerpoint(file_path: str) -> DocumentContent:
        """Extract content from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            text = ""
            slides = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {slide_num + 1} ---\n"
                slide_content = {"slide_number": slide_num + 1, "text": ""}
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + '\n'
                        slide_content["text"] += shape.text + '\n'
                
                text += slide_text
                slides.append(slide_content)
            
            metadata = {
                "num_slides": len(prs.slides),
                "title": prs.core_properties.title or '',
                "author": prs.core_properties.author or '',
                "subject": prs.core_properties.subject or '',
            }
            
            return DocumentContent(text=text, metadata=metadata, slides=slides)
            
        except Exception as e:
            return DocumentContent(text=f"Error extracting PowerPoint: {str(e)}", metadata={})
    
    @staticmethod
    def extract_word(file_path: str) -> DocumentContent:
        """Extract content from Word documents"""
        try:
            doc = Document(file_path)
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + '\n'
            
            # Extract table content
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + '\t'
                    text += '\n'
            
            metadata = {
                "num_paragraphs": len(doc.paragraphs),
                "num_tables": len(doc.tables),
                "title": doc.core_properties.title or '',
                "author": doc.core_properties.author or '',
                "subject": doc.core_properties.subject or '',
            }
            
            return DocumentContent(text=text, metadata=metadata)
            
        except Exception as e:
            return DocumentContent(text=f"Error extracting Word: {str(e)}", metadata={})
    
    @staticmethod
    def extract_csv(file_path: str) -> DocumentContent:
        """Extract content from CSV files"""
        try:
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            
            metadata = {
                "num_rows": len(df),
                "num_columns": len(df.columns),
                "column_names": df.columns.tolist(),
                "data_types": df.dtypes.astype(str).to_dict()
            }
            
            return DocumentContent(text=text, metadata=metadata)
            
        except Exception as e:
            return DocumentContent(text=f"Error extracting CSV: {str(e)}", metadata={})
    
    @staticmethod
    def extract_json(file_path: str) -> DocumentContent:
        """Extract content from JSON files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            text = json.dumps(data, indent=2, ensure_ascii=False)
            
            metadata = {
                "type": type(data).__name__,
                "size": len(str(data)) if isinstance(data, (dict, list)) else 0
            }
            
            return DocumentContent(text=text, metadata=metadata)
            
        except Exception as e:
            return DocumentContent(text=f"Error extracting JSON: {str(e)}", metadata={})
    
    @staticmethod
    def extract_xml(file_path: str) -> DocumentContent:
        """Extract content from XML files"""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            def extract_text(element):
                text = element.text or ""
                for child in element:
                    text += extract_text(child)
                return text
            
            text = extract_text(root)
            
            metadata = {
                "root_tag": root.tag,
                "num_elements": len(list(root.iter()))
            }
            
            return DocumentContent(text=text, metadata=metadata)
            
        except Exception as e:
            return DocumentContent(text=f"Error extracting XML: {str(e)}", metadata={})

class SemanticSearchEngine:
    """Handles semantic search using sentence transformers"""
    
    def __init__(self, model_name: str = "all-MiniLM-L6-v2"):
        self.model = SentenceTransformer(model_name)
        self.embeddings_cache = {}
        self.file_contents = {}
        
    def get_embedding(self, text: str) -> np.ndarray:
        """Get embedding for text"""
        text_hash = hashlib.md5(text.encode()).hexdigest()
        if text_hash not in self.embeddings_cache:
            self.embeddings_cache[text_hash] = self.model.encode([text])[0]
        return self.embeddings_cache[text_hash]
    
    def add_document(self, file_path: str, content: str, metadata: Dict[str, Any]):
        """Add document to search index"""
        self.file_contents[file_path] = {
            "content": content,
            "metadata": metadata,
            "embedding": self.get_embedding(content)
        }
    
    def search(self, query: str, top_k: int = 10, similarity_threshold: float = 0.1) -> List[Dict[str, Any]]:
        """Perform semantic search"""
        query_embedding = self.get_embedding(query)
        results = []
        
        for file_path, doc_data in self.file_contents.items():
            similarity = cosine_similarity(
                [query_embedding], 
                [doc_data["embedding"]]
            )[0][0]
            
            if similarity >= similarity_threshold:
                results.append({
                    "file_path": file_path,
                    "similarity": float(similarity),
                    "content": doc_data["content"],
                    "metadata": doc_data["metadata"]
                })
        
        results.sort(key=lambda x: x["similarity"], reverse=True)
        return results[:top_k]
    
    def save_index(self, file_path: str):
        """Save search index to file"""
        with open(file_path, 'wb') as f:
            pickle.dump({
                "embeddings_cache": self.embeddings_cache,
                "file_contents": self.file_contents
            }, f)
    
    def load_index(self, file_path: str):
        """Load search index from file"""
        try:
            with open(file_path, 'rb') as f:
                data = pickle.load(f)
                self.embeddings_cache = data.get("embeddings_cache", {})
                self.file_contents = data.get("file_contents", {})
        except Exception as e:
            print(f"Warning: Could not load search index: {e}")

# Create MCP instance  
mcp = FastMCP("Semantic Filesystem Server 🧠")

class FilesystemServer:
    def __init__(self, root_path: str, custom_ignore_patterns: Optional[List[str]] = None):
        """Initialize the semantic filesystem server"""
        self.root_path = os.path.abspath(root_path)
        if not os.path.exists(self.root_path):
            raise ValueError(f"Directory does not exist: {self.root_path}")
        
        # Initialize ignore patterns
        patterns = DEFAULT_IGNORE_PATTERNS.copy()
        if custom_ignore_patterns:
            patterns.extend(custom_ignore_patterns)
        
        # Load .gitignore if exists
        gitignore_path = os.path.join(self.root_path, ".gitignore")
        if os.path.exists(gitignore_path):
            with open(gitignore_path, "r") as f:
                patterns.extend(line.strip() for line in f if line.strip() and not line.startswith("#"))
        
        self.ignore_patterns = pathspec.PathSpec.from_lines("gitwildmatch", patterns)
        
        # Initialize document extractor and semantic search
        self.extractor = DocumentExtractor()
        self.search_engine = SemanticSearchEngine()
        
        # Load existing search index
        index_path = os.path.join(self.root_path, ".semantic_index.pkl")
        self.search_engine.load_index(index_path)
        
        # Register tools, resources, and prompts
        self._register_tools()
        self._register_resources()
        self._register_prompts()
        
        # Build initial index
        self._build_search_index()
    
    def _is_safe_path(self, path: str) -> bool:
        """Check if a path is safe to access"""
        abs_path = os.path.abspath(os.path.join(self.root_path, path))
        return abs_path.startswith(self.root_path)
    
    def _is_ignored(self, path: str) -> bool:
        """Check if path matches ignore patterns"""
        rel_path = os.path.relpath(path, self.root_path)
        return self.ignore_patterns.match_file(rel_path)
    
    def _get_file_metadata(self, file_path: str) -> Optional[FileMetadata]:
        """Get comprehensive file metadata"""
        try:
            if not os.path.exists(file_path) or not os.path.isfile(file_path):
                return None
            
            stat = os.stat(file_path)
            mime_type, encoding = mimetypes.guess_type(file_path)
            mime_type = mime_type or "application/octet-stream"
            encoding = encoding or "utf-8"
            
            # Determine content type and extractability
            content_type = self._determine_content_type(file_path, mime_type)
            extractable = self._is_extractable(file_path, mime_type)
            
            # Calculate file hash
            with open(file_path, "rb") as f:
                file_hash = hashlib.md5(f.read()).hexdigest()
            
            return FileMetadata(
                path=os.path.relpath(file_path, self.root_path),
                size=stat.st_size,
                modified=datetime.fromtimestamp(stat.st_mtime),
                mime_type=mime_type,
                encoding=encoding,
                hash=file_hash,
                content_type=content_type,
                extractable=extractable
            )
        except Exception:
            return None
    
    def _determine_content_type(self, file_path: str, mime_type: str) -> str:
        """Determine high-level content type"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.pdf']:
            return 'document'
        elif ext in ['.docx', '.doc']:
            return 'document'
        elif ext in ['.xlsx', '.xls', '.csv']:
            return 'spreadsheet'
        elif ext in ['.pptx', '.ppt']:
            return 'presentation'
        elif ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.yaml', '.yml']:
            return 'text'
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
            return 'image'
        elif ext in ['.mp3', '.wav', '.flac', '.aac']:
            return 'audio'
        elif ext in ['.mp4', '.avi', '.mov', '.wmv', '.flv']:
            return 'video'
        else:
            return 'unknown'
    
    def _is_extractable(self, file_path: str, mime_type: str) -> bool:
        """Check if content can be extracted from file"""
        ext = os.path.splitext(file_path)[1].lower()
        extractable_extensions = [
            '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.csv',
            '.pptx', '.ppt', '.txt', '.md', '.py', '.js',
            '.html', '.css', '.json', '.xml', '.yaml', '.yml'
        ]
        return ext in extractable_extensions
    
    def _extract_content(self, file_path: str) -> Optional[DocumentContent]:
        """Extract content from various file types"""
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == '.pdf':
                return self.extractor.extract_pdf(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self.extractor.extract_excel(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self.extractor.extract_powerpoint(file_path)
            elif ext in ['.docx', '.doc']:
                return self.extractor.extract_word(file_path)
            elif ext == '.csv':
                return self.extractor.extract_csv(file_path)
            elif ext == '.json':
                return self.extractor.extract_json(file_path)
            elif ext in ['.xml']:
                return self.extractor.extract_xml(file_path)
            elif ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.yaml', '.yml']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return DocumentContent(text=content, metadata={'type': 'text'})
            else:
                return None
        except Exception as e:
            return DocumentContent(text=f"Error extracting content: {str(e)}", metadata={})
    
    def _build_search_index(self):
        """Build semantic search index for all files"""
        print("Building semantic search index...")
        
        for root, dirs, files in os.walk(self.root_path):
            dirs[:] = [d for d in dirs if not self._is_ignored(os.path.join(root, d))]
            
            for file in files:
                full_path = os.path.join(root, file)
                
                if self._is_ignored(full_path):
                    continue
                
                metadata = self._get_file_metadata(full_path)
                if not metadata or not metadata.extractable:
                    continue
                
                content = self._extract_content(full_path)
                if content and content.text:
                    rel_path = os.path.relpath(full_path, self.root_path)
                    self.search_engine.add_document(
                        rel_path,
                        content.text,
                        {
                            "file_metadata": metadata.__dict__,
                            "content_metadata": content.metadata
                        }
                    )
        
        # Save index
        index_path = os.path.join(self.root_path, ".semantic_index.pkl")
        self.search_engine.save_index(index_path)
        print("Semantic search index built successfully!")
        
    def _extract_snippet(self, text: str, query: str, context_size: int = 100) -> str:
        """Extract a snippet around the query match"""
        query_lower = query.lower()
        text_lower = text.lower()
        
        match_pos = text_lower.find(query_lower)
        if match_pos == -1:
            return text[:200] + "..." if len(text) > 200 else text
        
        start = max(0, match_pos - context_size)
        end = min(len(text), match_pos + len(query) + context_size)
        
        snippet = text[start:end]
        if start > 0:
            snippet = "..." + snippet
        if end < len(text):
            snippet = snippet + "..."
        
        return snippet
    
    def _register_tools(self):
        """Register all tools"""
    
        @mcp.tool()
        def semantic_search(
            query: str,
            top_k: int = 10,
            similarity_threshold: float = 0.5,
            file_types: Optional[List[str]] = None,
            include_content: bool = True
        ) -> Dict[str, Any]:
            """
            Perform semantic search across all indexed documents.
            
            Args:
                query: Natural language search query
                top_k: Maximum number of results to return
                similarity_threshold: Minimum similarity score (0-1)
                file_types: Filter by file types (e.g., ['pdf', 'docx'])
                include_content: Whether to include content snippets
            
            Returns:
                Dictionary with semantic search results
            """
            results = self.search_engine.search(query, top_k, similarity_threshold)
            
            # Filter by file types if specified
            if file_types:
                file_types = [ft.lower() for ft in file_types]
                results = [r for r in results if any(r['file_path'].lower().endswith(ft) for ft in file_types)]
            
            # Format results
            formatted_results = []
            for result in results:
                formatted_result = {
                    "file_path": result["file_path"],
                    "similarity_score": result["similarity"],
                    "file_metadata": result["metadata"].get("file_metadata", {}),
                    "content_metadata": result["metadata"].get("content_metadata", {})
                }
                
                if include_content:
                    # Extract relevant snippets
                    content = result["content"]
                    sentences = content.split('.')
                    query_words = query.lower().split()
                    
                    relevant_snippets = []
                    for sentence in sentences[:50]:  # Limit to first 50 sentences
                        if any(word in sentence.lower() for word in query_words):
                            relevant_snippets.append(sentence.strip())
                    
                    formatted_result["content_snippets"] = relevant_snippets[:5]  # Top 5 snippets
                
                formatted_results.append(formatted_result)
            
            return {
                "query": query,
                "results": formatted_results,
                "total_results": len(formatted_results),
                "search_type": "semantic"
            }
        
        @mcp.tool()
        def extract_document_content(
            path: str,
            include_metadata: bool = True,
            include_structure: bool = False
        ) -> Dict[str, Any]:
            """
            Extract full content from a document.
            
            Args:
                path: Path to the document
                include_metadata: Whether to include document metadata
                include_structure: Whether to include structural information
            
            Returns:
                Dictionary with extracted content and metadata
            """
            if not self._is_safe_path(path):
                return {"error": "Path is outside root directory"}
            
            full_path = os.path.join(self.root_path, path)
            
            if not os.path.exists(full_path):
                return {"error": "File not found"}
            
            metadata = self._get_file_metadata(full_path)
            if not metadata or not metadata.extractable:
                return {"error": "File type not supported for content extraction"}
            
            content = self._extract_content(full_path)
            if not content:
                return {"error": "Failed to extract content"}
            
            result = {
                "file_path": path,
                "content": content.text,
                "extraction_successful": True
            }
            
            if include_metadata:
                result["file_metadata"] = metadata.__dict__
                result["content_metadata"] = content.metadata
            
            if include_structure:
                if content.pages:
                    result["pages"] = content.pages
                if content.sheets:
                    result["sheets"] = content.sheets
                if content.slides:
                    result["slides"] = content.slides
            
            return result
        
        @mcp.tool()
        def rebuild_search_index(
            incremental: bool = False
        ) -> Dict[str, Any]:
            """
            Rebuild the semantic search index.
            
            Args:
                incremental: Whether to do incremental update (only new/changed files)
            
            Returns:
                Dictionary with rebuild status
            """
            try:
                if not incremental:
                    self.search_engine.file_contents.clear()
                    self.search_engine.embeddings_cache.clear()
                
                self._build_search_index()
                
                return {
                    "success": True,
                    "message": "Search index rebuilt successfully",
                    "indexed_files": len(self.search_engine.file_contents)
                }
            except Exception as e:
                return {
                    "success": False,
                    "error": f"Failed to rebuild index: {str(e)}"
                }
    
        @mcp.tool()
        def get_file_summary(
            path: str,
            summary_type: str = "auto"
        ) -> Dict[str, Any]:
            """
            Get AI-generated summary of a file's content.
            
            Args:
                path: Path to the file
                summary_type: Type of summary (auto, technical, business, academic)
            
            Returns:
                Dictionary with file summary and analysis
            """
            if not self._is_safe_path(path):
                return {"error": "Path is outside root directory"}
            
            full_path = os.path.join(self.root_path, path)
            
            if not os.path.exists(full_path):
                return {"error": "File not found"}
            
            metadata = self._get_file_metadata(full_path)
            if not metadata or not metadata.extractable:
                return {"error": "File type not supported for content extraction"}
            
            content = self._extract_content(full_path)
            if not content:
                return {"error": "Failed to extract content"}
            
            # Basic content analysis
            word_count = len(content.text.split())
            char_count = len(content.text)
            
            # Extract key information based on content type
            summary = {
                "file_path": path,
                "content_type": metadata.content_type,
                "word_count": word_count,
                "character_count": char_count,
                "file_metadata": metadata.__dict__,
                "content_metadata": content.metadata
            }
            
            # Add type-specific analysis
            if metadata.content_type == "document":
                summary["document_analysis"] = {
                    "estimated_reading_time": f"{max(1, word_count // 200)} minutes",
                    "content_preview": content.text[:500] + "..." if len(content.text) > 500 else content.text
                }
            elif metadata.content_type == "spreadsheet":
                summary["spreadsheet_analysis"] = {
                    "data_summary": content.metadata,
                    "data_preview": content.text[:1000] + "..." if len(content.text) > 1000 else content.text
                }
            elif metadata.content_type == "presentation":
                summary["presentation_analysis"] = {
                    "slide_count": content.metadata.get("num_slides", 0),
                    "content_preview": content.text[:500] + "..." if len(content.text) > 500 else content.text
                }
            
            return summary
    
        # Basic file operations
        @mcp.tool()
        def list_files(
            path: str = "",
            include_hidden: bool = False,
            pattern: Optional[str] = None,
            content_type: Optional[str] = None
        ) -> Dict[str, Any]:
            """
            List files and directories in the given path.
            
            Args:
                path: Relative path from root directory
                include_hidden: Whether to include hidden files
                pattern: Glob pattern to filter files
                content_type: Filter by content type (document, spreadsheet, text, etc.)
            
            Returns:
                Dictionary with file listing information
            """
            if not self._is_safe_path(path):
                return {"error": "Path is outside root directory"}
            
            full_path = os.path.join(self.root_path, path)
            
            if not os.path.exists(full_path):
                return {"error": "Directory not found"}
            
            if not os.path.isdir(full_path):
                return {"error": "Path is not a directory"}
            
            files = []
            directories = []
            
            try:
                for item in os.listdir(full_path):
                    item_path = os.path.join(full_path, item)
                    
                    if not include_hidden and item.startswith('.'):
                        continue
                    
                    if self._is_ignored(item_path):
                        continue
                    
                    if pattern and not fnmatch.fnmatch(item, pattern):
                        continue
                    
                    if os.path.isdir(item_path):
                        directories.append({
                            "name": item,
                            "path": os.path.relpath(item_path, self.root_path),
                            "type": "directory"
                        })
                    else:
                        metadata = self._get_file_metadata(item_path)
                        if metadata:
                            if content_type and metadata.content_type != content_type:
                                continue
                            
                            files.append({
                                "name": item,
                                "path": metadata.path,
                                "size": metadata.size,
                                "modified": metadata.modified.isoformat(),
                                "mime_type": metadata.mime_type,
                                "content_type": metadata.content_type,
                                "extractable": metadata.extractable,
                                "hash": metadata.hash
                            })
                
                return {
                    "path": path,
                    "files": files,
                    "directories": directories,
                    "total_files": len(files),
                    "total_directories": len(directories)
                }
            except Exception as e:
                return {"error": f"Failed to list directory: {str(e)}"}
    
        @mcp.tool()
        def read_file(
            path: str,
            encoding: str = "utf-8",
            max_size: int = 10 * 1024 * 1024  # 10MB limit
        ) -> Dict[str, Any]:
            """
            Read the contents of a file.
            
            Args:
                path: Relative path from root directory
                encoding: Text encoding to use
                max_size: Maximum file size to read (bytes)
            
            Returns:
                Dictionary with file contents and metadata
            """
            if not self._is_safe_path(path):
                return {"error": "Path is outside root directory"}
            
            full_path = os.path.join(self.root_path, path)
            
            if not os.path.exists(full_path):
                return {"error": "File not found"}
            
            if not os.path.isfile(full_path):
                return {"error": "Path is not a file"}
            
            try:
                file_size = os.path.getsize(full_path)
                if file_size > max_size:
                    return {"error": f"File too large: {file_size} bytes > {max_size} bytes"}
                
                metadata = self._get_file_metadata(full_path)
                
                # Try to read as text first
                try:
                    with open(full_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    
                    return {
                        "path": path,
                        "content": content,
                        "size": file_size,
                        "encoding": encoding,
                        "metadata": metadata.__dict__ if metadata else None,
                        "type": "text"
                    }
                except UnicodeDecodeError:
                    # If text reading fails, return binary info
                    return {
                        "path": path,
                        "error": "Binary file - use extract_document_content for supported formats",
                        "size": file_size,
                        "metadata": metadata.__dict__ if metadata else None,
                        "type": "binary"
                    }
                    
            except Exception as e:
                return {"error": f"Failed to read file: {str(e)}"}
    
        @mcp.tool()
        def write_file(
            path: str,
            content: str,
            encoding: str = "utf-8",
            create_dirs: bool = True
        ) -> Dict[str, Any]:
            """
            Write content to a file.
            
            Args:
                path: Relative path from root directory
                content: Content to write
                encoding: Text encoding to use
                create_dirs: Whether to create parent directories
            
            Returns:
                Dictionary with write operation status
            """
            if not self._is_safe_path(path):
                return {"error": "Path is outside root directory"}
            
            full_path = os.path.join(self.root_path, path)
            
            try:
                if create_dirs:
                    os.makedirs(os.path.dirname(full_path), exist_ok=True)
                
                with open(full_path, 'w', encoding=encoding) as f:
                    f.write(content)
                
                # Update search index if file is extractable
                metadata = self._get_file_metadata(full_path)
                if metadata and metadata.extractable:
                    extracted_content = self._extract_content(full_path)
                    if extracted_content:
                        self.search_engine.add_document(
                            path,
                            extracted_content.text,
                            {
                                "file_metadata": metadata.__dict__,
                                "content_metadata": extracted_content.metadata
                            }
                        )
                
                return {
                    "success": True,
                    "path": path,
                    "bytes_written": len(content.encode(encoding)),
                    "message": "File written successfully"
                }
                
            except Exception as e:
                return {"error": f"Failed to write file: {str(e)}"}
    
        @mcp.tool()
        def search_files(
            query: str,
            search_type: str = "hybrid",
            path: str = "",
            file_extensions: Optional[List[str]] = None,
            modified_after: Optional[str] = None,
            modified_before: Optional[str] = None,
            max_results: int = 20
        ) -> Dict[str, Any]:
            """
            Search for files using various methods.
            
            Args:
                query: Search query (filename, content, or semantic)
                search_type: Type of search (filename, content, semantic, hybrid)
                path: Limit search to specific path
                file_extensions: Filter by file extensions
                modified_after: ISO format date string
                modified_before: ISO format date string
                max_results: Maximum number of results
            
            Returns:
                Dictionary with search results
            """
            if not self._is_safe_path(path):
                return {"error": "Path is outside root directory"}
            
            search_path = os.path.join(self.root_path, path)
            results = []
            
            try:
                # Parse date filters
                after_date = None
                before_date = None
                if modified_after:
                    after_date = datetime.fromisoformat(modified_after)
                if modified_before:
                    before_date = datetime.fromisoformat(modified_before)
                
                if search_type in ["filename", "hybrid"]:
                    # Filename search
                    for root, dirs, files in os.walk(search_path):
                        dirs[:] = [d for d in dirs if not self._is_ignored(os.path.join(root, d))]
                        
                        for file in files:
                            full_path = os.path.join(root, file)
                            
                            if self._is_ignored(full_path):
                                continue
                            
                            # Apply filters
                            if file_extensions and not any(file.lower().endswith(ext.lower()) for ext in file_extensions):
                                continue
                            
                            metadata = self._get_file_metadata(full_path)
                            if not metadata:
                                continue
                            
                            if after_date and metadata.modified < after_date:
                                continue
                            if before_date and metadata.modified > before_date:
                                continue
                            
                            # Check if query matches filename
                            if query.lower() in file.lower():
                                results.append({
                                    "file_path": metadata.path,
                                    "match_type": "filename",
                                    "match_score": 1.0,
                                    "metadata": metadata.__dict__
                                })
                
                if search_type in ["content", "hybrid"]:
                    # Content search (simple text matching)
                    for root, dirs, files in os.walk(search_path):
                        dirs[:] = [d for d in dirs if not self._is_ignored(os.path.join(root, d))]
                        
                        for file in files:
                            full_path = os.path.join(root, file)
                            
                            if self._is_ignored(full_path):
                                continue
                            
                            metadata = self._get_file_metadata(full_path)
                            if not metadata or not metadata.extractable:
                                continue
                            
                            # Apply filters
                            if file_extensions and not any(file.lower().endswith(ext.lower()) for ext in file_extensions):
                                continue
                            
                            if after_date and metadata.modified < after_date:
                                continue
                            if before_date and metadata.modified > before_date:
                                continue
                            
                            # Extract and search content
                            content = self._extract_content(full_path)
                            if content and query.lower() in content.text.lower():
                                results.append({
                                    "file_path": metadata.path,
                                    "match_type": "content",
                                    "match_score": 0.8,
                                    "metadata": metadata.__dict__,
                                    "content_snippet": self._extract_snippet(content.text, query)
                                })
                
                if search_type in ["semantic", "hybrid"]:
                    # Semantic search
                    semantic_results = self.search_engine.search(query, max_results, 0.1)
                    
                    for result in semantic_results:
                        file_path = result["file_path"]
                        full_path = os.path.join(self.root_path, file_path)
                        
                        # Apply filters
                        if file_extensions and not any(file_path.lower().endswith(ext.lower()) for ext in file_extensions):
                            continue
                        
                        metadata = self._get_file_metadata(full_path)
                        if not metadata:
                            continue
                        
                        if after_date and metadata.modified < after_date:
                            continue
                        if before_date and metadata.modified > before_date:
                            continue
                        
                        results.append({
                            "file_path": file_path,
                            "match_type": "semantic",
                            "match_score": result["similarity"],
                            "metadata": metadata.__dict__,
                            "content_snippet": self._extract_snippet(result["content"], query)
                        })
                
                # Remove duplicates and sort by score
                seen_files = set()
                unique_results = []
                for result in results:
                    if result["file_path"] not in seen_files:
                        seen_files.add(result["file_path"])
                        unique_results.append(result)
                
                unique_results.sort(key=lambda x: x["match_score"], reverse=True)
                
                return {
                    "query": query,
                    "search_type": search_type,
                    "results": unique_results[:max_results],
                    "total_results": len(unique_results)
                }
                
            except Exception as e:
                return {"error": f"Search failed: {str(e)}"}
    
        @mcp.tool()
        def get_file_info(
            path: str,
            include_content_analysis: bool = False
        ) -> Dict[str, Any]:
            """
            Get detailed information about a file.
            
            Args:
                path: Relative path from root directory
                include_content_analysis: Whether to analyze content
            
            Returns:
                Dictionary with comprehensive file information
            """
            if not self._is_safe_path(path):
                return {"error": "Path is outside root directory"}
            
            full_path = os.path.join(self.root_path, path)
            
            if not os.path.exists(full_path):
                return {"error": "File not found"}
            
            metadata = self._get_file_metadata(full_path)
            if not metadata:
                return {"error": "Could not get file metadata"}
            
            result = {
                "path": path,
                "metadata": metadata.__dict__,
                "absolute_path": full_path
            }
            
            if include_content_analysis and metadata.extractable:
                content = self._extract_content(full_path)
                if content:
                    result["content_analysis"] = {
                        "word_count": len(content.text.split()),
                        "character_count": len(content.text),
                        "content_metadata": content.metadata,
                        "content_preview": content.text[:500] + "..." if len(content.text) > 500 else content.text
                    }
            
            return result

    def _register_resources(self):
        """Register MCP resources"""
        
        @mcp.resource("file://{path}")
        def get_file_resource(path: str) -> str:
            """Get file content as a resource"""
            if not self._is_safe_path(path):
                return "Error: Path is outside root directory"
            
            full_path = os.path.join(self.root_path, path)
            
            if not os.path.exists(full_path):
                return "Error: File not found"
            
            metadata = self._get_file_metadata(full_path)
            if not metadata:
                return "Error: Could not get file metadata"
            
            if metadata.extractable:
                content = self._extract_content(full_path)
                if content:
                    return content.text
            
            # Try to read as plain text
            try:
                with open(full_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except:
                return "Error: Could not read file content"

    def _register_prompts(self):
        """Register MCP prompts"""
        
        @mcp.prompt()
        def analyze_codebase(
            query: str = "Analyze the codebase structure and provide insights",
            focus_area: str = "general"
        ) -> str:
            """
            Analyze the codebase and provide insights.
            
            Args:
                query: Specific analysis question
                focus_area: Area to focus on (general, security, performance, architecture)
            """
            
            # Get code files
            code_files = []
            for root, dirs, files in os.walk(self.root_path):
                dirs[:] = [d for d in dirs if not self._is_ignored(os.path.join(root, d))]
                
                for file in files:
                    if file.endswith(('.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.go', '.rs')):
                        full_path = os.path.join(root, file)
                        if not self._is_ignored(full_path):
                            code_files.append(os.path.relpath(full_path, self.root_path))
            
            prompt = f"""
            You are analyzing a codebase with {len(code_files)} code files.
            
            Query: {query}
            Focus Area: {focus_area}
            
            Key files identified:
            {chr(10).join(code_files[:20])}  # Show first 20 files
            {"... and more" if len(code_files) > 20 else ""}
            
            Please provide insights about:
            1. Code structure and organization
            2. Key components and their relationships
            3. Potential areas for improvement
            4. Security considerations (if focus_area is security)
            5. Performance considerations (if focus_area is performance)
            6. Architecture patterns used (if focus_area is architecture)
            
            Use the semantic_search and extract_document_content tools to analyze specific files.
            """
            
            return prompt
        
        @mcp.prompt()
        def document_qa(
            question: str,
            document_types: str = "all"
        ) -> str:
            """
            Answer questions about documents in the filesystem.
            
            Args:
                question: Question about the documents
                document_types: Types of documents to consider (all, pdf, docx, etc.)
            """
            
            prompt = f"""
            You are helping answer questions about documents in the filesystem.
            
            Question: {question}
            Document Types: {document_types}
            
            Use the semantic_search tool to find relevant documents, then use extract_document_content 
            to get the full content of relevant documents to answer the question.
            
            Focus on providing accurate, well-sourced answers based on the document content.
            """
            
            return prompt

def main():
    """Main entry point"""
    import argparse
    
    parser = argparse.ArgumentParser(description="FastMCP Semantic Filesystem Server")
    parser.add_argument("root_path", help="Root directory to serve")
    parser.add_argument("--ignore", nargs="*", help="Additional ignore patterns")
    parser.add_argument("--rebuild-index", action="store_true", help="Rebuild search index on startup")
    
    args = parser.parse_args()
    
    # Validate root path
    if not os.path.exists(args.root_path):
        print(f"Error: Directory does not exist: {args.root_path}")
        sys.exit(1)
    
    try:        
        # Initialize the server
        server = FilesystemServer(args.root_path, args.ignore)

        if args.rebuild_index:
            server._build_search_index()
            
        print(f"Starting FastMCP Semantic Filesystem Server on {args.root_path}")
        mcp.run()
        
    except KeyboardInterrupt:
        print("\nShutting down server...")
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()